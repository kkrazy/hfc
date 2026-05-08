#!/usr/bin/env python3
"""Generate hfc_design.pptx — design + components deck for the lite serving layer.

Run: python3 scripts/gen_design_ppt.py
Output: out/hfc_design.pptx
"""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# --- palette (calm, professional) -------------------------------------------- #
NAVY = RGBColor(0x1B, 0x36, 0x5D)
TEAL = RGBColor(0x2E, 0x86, 0x8B)
GREY = RGBColor(0x55, 0x5B, 0x6E)
LIGHT = RGBColor(0xEE, 0xF1, 0xF4)
RED = RGBColor(0xB0, 0x2A, 0x2A)
GREEN = RGBColor(0x2A, 0x7A, 0x2A)
ORANGE = RGBColor(0xC8, 0x6A, 0x1E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x12, 0x14, 0x18)


def add_title_only(prs, title, subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(s, WHITE)
    t = add_text(s, 0.5, 0.4, 12.3, 0.9, title, size=30, bold=True, color=NAVY)
    if subtitle:
        add_text(s, 0.5, 1.25, 12.3, 0.5, subtitle, size=15, color=GREY, italic=True)
    return s


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=BLACK,
             italic=False, align=None, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=14, color=BLACK, font="Calibri"):
    """items: list of str OR (str, level) tuples for indent."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            text, lvl = it
        else:
            text, lvl = it, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        bullet = "• " if lvl == 0 else "– "
        r = p.add_run()
        r.text = bullet + text
        r.font.name = font
        r.font.size = Pt(size - lvl)
        r.font.color.rgb = color
    return tb


def add_box(slide, x, y, w, h, fill=LIGHT, border=NAVY, text=None,
            size=12, bold=False, text_color=BLACK):
    from pptx.shapes.autoshape import Shape
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border
    sh.line.width = Pt(1.25)
    if text:
        sh.text_frame.margin_left = Emu(50000)
        sh.text_frame.margin_right = Emu(50000)
        sh.text_frame.margin_top = Emu(30000)
        sh.text_frame.margin_bottom = Emu(30000)
        sh.text_frame.word_wrap = True
        p = sh.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = text_color
        r.font.name = "Calibri"
    return sh


def add_arrow(slide, x1, y1, x2, y2, color=NAVY, weight=2):
    from pptx.enum.shapes import MSO_CONNECTOR
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1),
                                   Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(weight)
    # Arrow head on end
    line = c.line
    from pptx.oxml.ns import qn
    lnXml = line._get_or_add_ln()
    tail = lnXml.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    lnXml.append(tail)
    return c


def add_footer(slide, page_n, total):
    add_text(slide, 0.4, 7.0, 12.5, 0.3,
             f"HFC — DRAM-Centric Inference  ·  {page_n}/{total}",
             size=9, color=GREY)


# ============================================================================ #
#  Slide builders
# ============================================================================ #


def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, NAVY)
    # Title block
    add_text(s, 0.7, 2.2, 12.0, 1.2, "HFC Serving Layer",
             size=48, bold=True, color=WHITE)
    add_text(s, 0.7, 3.3, 12.0, 0.7,
             "DRAM-Centric Inference on Ascend NPU",
             size=24, color=LIGHT)
    add_text(s, 0.7, 4.1, 12.0, 0.5,
             "Component design for collaborative implementation",
             size=16, italic=True, color=LIGHT)
    add_text(s, 0.7, 5.3, 12.0, 0.4,
             "Lite layer over HF transformers  •  Continuous batching  •  "
             "DRAM-resident KV  •  Streamed submission",
             size=12, color=LIGHT)
    add_text(s, 0.7, 7.0, 12.0, 0.3,
             "2026-05-07  ·  v0",
             size=10, color=LIGHT)
    return s


def slide_why(prs, n, total):
    s = add_title_only(prs, "Why build this", "HF + a thin scheduler beats vLLM for our workload")
    add_bullets(s, 0.5, 1.7, 12.0, 5.0, [
        "Hardware: 8× Ascend 910 (16 logical NPUs), 64 GB HBM/device, 2 TB host DRAM",
        ("HBM scarcity is not our bottleneck — DRAM headroom inverts vLLM's design constraint", 1),
        "HF transformers does NOT have real continuous batching, real PP, or DRAM offload",
        ("`device_map=\"auto\"` runs sequentially — capacity, not throughput", 1),
        "vLLM has these but uses paged attention (designed for HBM scarcity, not our problem)",
        ("PagedAttention adds complexity we don't need on a DRAM-rich box", 1),
        "Goal: thin layer over HF that exploits 2 TB DRAM and serves requests immediately",
        ("Custom `transformers.Cache` subclass = clean integration point", 1),
        ("FX-rewriter pipeline (already built) = the offload mechanism", 1),
        ("Continuous batching + streamed submission = the throughput levers", 1),
    ], size=14)
    add_footer(s, n, total)


def slide_hardware(prs, n, total):
    s = add_title_only(prs, "Hardware target — measured ceiling",
                       "Practical peak ≠ rated peak; design around what we measured")
    # Box: practical numbers
    add_box(s, 0.5, 1.7, 6.0, 0.6, fill=LIGHT, border=NAVY,
            text="Ascend 910 9392  ·  8 cards × 2 logical NPUs = 16",
            size=14, bold=True)
    add_bullets(s, 0.5, 2.5, 6.0, 4.5, [
        "64 GB HBM per logical NPU",
        "2 TB host DRAM (shared across box)",
        "Practical peak: ~320 TFLOPS sustained",
        ("(86% of rated 376 TFLOPS fp16)", 1),
        "Per-kernel device floor: ~6.8 µs",
        ("irreducible firmware launch cost", 1),
        "Per-call sync overhead: ~35 µs flat",
        ("relevant for compute-bound submissions", 1),
        "HBM ↔ DRAM bandwidth: ~32 GB/s pinned",
    ], size=13)

    add_box(s, 6.8, 1.7, 6.0, 0.6, fill=LIGHT, border=TEAL,
            text="Implications for design", size=14, bold=True)
    add_bullets(s, 6.8, 2.5, 6.0, 4.5, [
        "Plan capacity around 320 TF/NPU, not 376",
        "Avoid sync between forwards (loses 20-30%)",
        "Big batch is THE lever (257× throughput gap 1→4096)",
        "Per-launch tricks bounded by 6.8 µs floor",
        ("only matters for small kernels (dim ≤ 2048)", 1),
        "DRAM offload cost amortizes over many decode steps",
        ("KV touched every step ⇒ per-byte cost spread out", 1),
    ], size=13)
    add_footer(s, n, total)


def slide_findings_helps(prs, n, total):
    s = add_title_only(prs, "What helps — measured",
                       "Findings ranked by impact on Qwen3-scale workloads")
    rows = [
        ("Continuous batching",      "1× → 4096× = 257× throughput gap",          "huge",  GREEN),
        ("Streamed submission",      "+5–30% at compute-bound (no per-call sync)", "easy",  GREEN),
        ("NPUGraph (small kernels)", "1.2× at dim ≤ 2048; ~null at large dims",    "yes",   ORANGE),
        ("DRAM KV offload",          "Removes per-request HBM pressure",           "yes",   GREEN),
        ("FX rewriter for offload",  "Composes with above; one capture, replayable","yes",   GREEN),
    ]
    y = 1.7
    add_text(s, 0.5, y, 4.5, 0.4, "Technique", size=14, bold=True, color=NAVY)
    add_text(s, 5.0, y, 5.5, 0.4, "Effect",    size=14, bold=True, color=NAVY)
    add_text(s, 11.0, y, 1.8, 0.4, "Verdict",  size=14, bold=True, color=NAVY)
    y += 0.5
    for label, effect, verdict, color in rows:
        add_box(s, 0.5, y, 4.4, 0.55, fill=LIGHT, border=NAVY,
                text=label, size=12, bold=True)
        add_box(s, 5.0, y, 5.9, 0.55, fill=WHITE, border=GREY,
                text=effect, size=11)
        add_box(s, 11.0, y, 1.7, 0.55, fill=color, border=color,
                text=verdict, size=12, bold=True, text_color=WHITE)
        y += 0.65
    add_footer(s, n, total)


def slide_findings_skip(prs, n, total):
    s = add_title_only(prs, "What does NOT help — measured",
                       "Empirical non-results — don't burn cycles on these")
    rows = [
        ("Multi-stream concurrent submission",
         "Tested 8 streams: 1.0× (best), 0.22× (worst). AICores serialize matmul.",
         "skip"),
        ("make_fx graph replay",
         "2.3× SLOWER than eager — graph holds 64 live intermediates per call.",
         "skip"),
        ("Tensor parallelism",
         "Per-op rewriting + all-reduce-after-every-linear; too invasive for our targets.",
         "skip"),
        ("HF device_map=\"auto\" as PP",
         "Sequential execution; only 1 device active at a time. ≠ real PP.",
         "skip"),
    ]
    y = 1.7
    for label, effect, verdict in rows:
        add_box(s, 0.5, y, 4.4, 0.95, fill=LIGHT, border=RED,
                text=label, size=12, bold=True)
        add_box(s, 5.0, y, 5.9, 0.95, fill=WHITE, border=GREY,
                text=effect, size=11)
        add_box(s, 11.0, y, 1.7, 0.95, fill=RED, border=RED,
                text=verdict, size=14, bold=True, text_color=WHITE)
        y += 1.05
    add_footer(s, n, total)


def slide_topology(prs, n, total):
    s = add_title_only(prs, "Architecture — Flink-style operator topology",
                       "Async DAG; queues bound backpressure; one operator owns the NPU")

    # HTTP source
    add_box(s, 0.4, 1.6, 1.6, 0.5, fill=NAVY, border=NAVY,
            text="HTTP /generate", size=10, bold=True, text_color=WHITE)

    # Admission queue
    add_box(s, 2.2, 1.65, 1.0, 0.4, fill=LIGHT, border=GREY, text="Q", size=10, bold=True)

    # Tokenize
    add_box(s, 3.4, 1.55, 1.6, 0.6, fill=TEAL, border=TEAL,
            text="Tokenize\n(parallel)", size=10, bold=True, text_color=WHITE)
    add_box(s, 5.2, 1.65, 1.0, 0.4, fill=LIGHT, border=GREY, text="Q", size=10, bold=True)

    # Scheduler
    add_box(s, 6.4, 1.55, 1.8, 0.6, fill=TEAL, border=TEAL,
            text="Scheduler\n(prefill/decode)", size=10, bold=True, text_color=WHITE)
    add_box(s, 8.4, 1.65, 1.0, 0.4, fill=LIGHT, border=GREY, text="Q", size=10, bold=True)

    # Forward (NPU)
    add_box(s, 9.6, 1.4, 3.2, 0.95, fill=NAVY, border=NAVY,
            text="Forward operator\n(NPU, parallelism = 1)\nstreamed, NPUGraph, FX-rewritten",
            size=10, bold=True, text_color=WHITE)

    # KV pool (side)
    add_box(s, 9.7, 2.6, 3.0, 0.6, fill=ORANGE, border=ORANGE,
            text="DRAM KV pool (contiguous per request)",
            size=10, bold=True, text_color=WHITE)

    # Logits Q
    add_box(s, 10.5, 3.5, 1.4, 0.4, fill=LIGHT, border=GREY,
            text="logits Q (futures)", size=10, bold=True)

    # Sampler
    add_box(s, 8.6, 4.2, 1.7, 0.6, fill=TEAL, border=TEAL,
            text="Sampler\n(parallel)", size=10, bold=True, text_color=WHITE)

    # Detokenize
    add_box(s, 6.5, 4.2, 1.7, 0.6, fill=TEAL, border=TEAL,
            text="Detokenize", size=10, bold=True, text_color=WHITE)

    # Per-request streams
    add_box(s, 4.0, 4.2, 2.1, 0.6, fill=NAVY, border=NAVY,
            text="Per-request streams\n→ client",
            size=10, bold=True, text_color=WHITE)

    # Arrows: top row left to right
    add_arrow(s, 2.0, 1.85, 2.2, 1.85)
    add_arrow(s, 3.2, 1.85, 3.4, 1.85)
    add_arrow(s, 5.0, 1.85, 5.2, 1.85)
    add_arrow(s, 6.2, 1.85, 6.4, 1.85)
    add_arrow(s, 8.2, 1.85, 8.4, 1.85)
    add_arrow(s, 9.4, 1.85, 9.6, 1.85)
    # Forward to logits Q
    add_arrow(s, 11.2, 2.4, 11.2, 3.5)
    # KV pool ↔ Forward
    add_arrow(s, 11.2, 2.6, 11.2, 2.4, color=ORANGE)
    # Logits Q to sampler
    add_arrow(s, 10.5, 3.7, 10.3, 4.4)
    # Sampler to detokenize
    add_arrow(s, 8.6, 4.5, 8.2, 4.5)
    # Detokenize to streams
    add_arrow(s, 6.5, 4.5, 6.1, 4.5)

    # Caption / labels
    add_text(s, 0.4, 5.4, 12.5, 0.4,
             "Async / asyncio runtime. Bounded queues = backpressure. "
             "Forward op never syncs; sampler triggers logits→host transfer.",
             size=12, italic=True, color=GREY)

    # Phase-mapping callouts
    add_text(s, 0.4, 5.9, 12.5, 0.3,
             "Phases:  P0 single-req  ·  P1 continuous batching  ·  "
             "P2 DRAM KV  ·  P3 streamed submission  ·  P4 NPUGraph + FX offload  ·  P5 PP + EP",
             size=11, color=NAVY)
    add_footer(s, n, total)


def slide_two_components(prs, n, total):
    s = add_title_only(prs, "Two components, one contract",
                       "Split for parallel ownership; interface boundary keeps changes local")

    # Component A
    add_box(s, 0.4, 1.7, 6.1, 4.8, fill=LIGHT, border=NAVY,
            text="", size=10)
    add_text(s, 0.6, 1.85, 5.8, 0.4, "Component A — Model Runtime",
             size=18, bold=True, color=NAVY)
    add_text(s, 0.6, 2.3, 5.8, 0.3,
             "Owns the NPU. Below the request abstraction.",
             size=11, italic=True, color=GREY)
    add_bullets(s, 0.6, 2.7, 5.8, 3.6, [
        "HF model integration + custom Cache",
        "DRAM KV pool (contiguous per request)",
        "Eager → NPUGraph → FX-rewritten forward",
        "Offload backend (D2H/H2D pinned memory)",
        "Owner: torch_npu / FX / model internals",
        "Files: forward_runner.py, dram_cache.py,",
        ("kv_pool.py, model_load.py, npu_graph.py,", 1),
        ("fx_offload.py", 1),
        "Phases: 0 → 1 → 2 → 3 → 4 → 5",
    ], size=12)

    # Component B
    add_box(s, 6.7, 1.7, 6.1, 4.8, fill=LIGHT, border=TEAL,
            text="", size=10)
    add_text(s, 6.9, 1.85, 5.8, 0.4, "Component B — Serving Engine",
             size=18, bold=True, color=TEAL)
    add_text(s, 6.9, 2.3, 5.8, 0.3,
             "Owns the request lifecycle. Above the model abstraction.",
             size=11, italic=True, color=GREY)
    add_bullets(s, 6.9, 2.7, 5.8, 3.6, [
        "Async operator topology (asyncio + queues)",
        "Request state, sampling params",
        "Tokenize / sampler / detokenize operators",
        "Scheduler (continuous batching)",
        "FastAPI server + streaming output",
        "Owner: async Python / web / scheduling",
        "Files: topology.py, server.py, operators/*.py,",
        ("state/request.py", 1),
        "Phases: 0 → 1 → 3 (others mostly transparent)",
    ], size=12)

    # Contract
    add_box(s, 0.4, 6.0, 12.4, 0.7, fill=ORANGE, border=ORANGE,
            text="Contract: PrefillBatch / DecodeBatch / StepResult / admit_request / release_request",
            size=14, bold=True, text_color=WHITE)
    add_footer(s, n, total)


def slide_component_a(prs, n, total):
    s = add_title_only(prs, "Component A — Model Runtime",
                       "Below the ForwardRunner interface; nothing visible to serving layer")

    add_text(s, 0.5, 1.7, 6.0, 0.3, "Files & responsibility",
             size=14, bold=True, color=NAVY)
    add_bullets(s, 0.5, 2.05, 6.0, 4.7, [
        "hf_integration/forward_runner.py",
        ("public ForwardRunner class — the contract", 1),
        ("eager (P0-1) → NPUGraph (P4) → FX-rewritten (P4)", 1),
        "hf_integration/dram_cache.py",
        ("DramKVCache(transformers.DynamicCache)", 1),
        ("override update() to redirect to KV pool", 1),
        "hf_integration/model_load.py",
        ("apply transformers/torch compat patches", 1),
        ("set _attn_implementation, dense MoE patch", 1),
        "state/kv_pool.py",
        ("DRAM allocator: max_concurrent × max_seq KV", 1),
        ("pool keyed (req_id, layer, k|v)", 1),
        "hf_integration/optimizations/npu_graph.py",
        ("aclgraph capture/replay for decode steps", 1),
        "hf_integration/optimizations/fx_offload.py",
        ("integrate existing HFC FX rewriter", 1),
    ], size=11)

    # Right column: phase mapping
    add_text(s, 7.2, 1.7, 5.5, 0.3, "What changes by phase",
             size=14, bold=True, color=NAVY)
    add_box(s, 7.2, 2.05, 5.5, 0.6, fill=LIGHT, border=NAVY,
            text="P0  HF generate(), single-shot per request", size=11, bold=True)
    add_box(s, 7.2, 2.7, 5.5, 0.6, fill=LIGHT, border=NAVY,
            text="P1  manual prefill/decode loop, HF DynamicCache", size=11, bold=True)
    add_box(s, 7.2, 3.35, 5.5, 0.6, fill=LIGHT, border=NAVY,
            text="P2  DramKVCache + kv_pool.py", size=11, bold=True)
    add_box(s, 7.2, 4.0, 5.5, 0.6, fill=LIGHT, border=NAVY,
            text="P3  streamed submission, no per-step sync", size=11, bold=True)
    add_box(s, 7.2, 4.65, 5.5, 0.6, fill=LIGHT, border=ORANGE,
            text="P4  NPUGraph capture + FX offload pass", size=11, bold=True)
    add_box(s, 7.2, 5.3, 5.5, 0.6, fill=LIGHT, border=TEAL,
            text="P5  partition_pipeline + EP rewrite (multi-NPU)", size=11, bold=True)

    add_text(s, 7.2, 6.05, 5.5, 0.4,
             "Day 1: ForwardRunner wraps model.generate.\n"
             "Phase 1: real prefill+decode tick.",
             size=11, italic=True, color=GREY)

    add_footer(s, n, total)


def slide_component_b(prs, n, total):
    s = add_title_only(prs, "Component B — Serving Engine",
                       "Above the ForwardRunner interface; never touches torch_npu")

    add_text(s, 0.5, 1.7, 6.0, 0.3, "Files & responsibility",
             size=14, bold=True, color=TEAL)
    add_bullets(s, 0.5, 2.05, 6.0, 4.7, [
        "topology.py",
        ("wire the operator DAG, lifecycle", 1),
        ("error propagation, graceful shutdown", 1),
        "operators/admission.py",
        ("HTTP request → admission queue", 1),
        "operators/tokenize.py",
        ("HF tokenizer wrapper, parallel=N", 1),
        "operators/scheduler.py",
        ("continuous batching, prefill vs decode tick", 1),
        ("HBM-pressure-aware admission", 1),
        "operators/forward.py",
        ("thin wrapper around ForwardRunner.step", 1),
        "operators/sampler.py",
        ("temperature/top-p, vectorized", 1),
        "operators/detokenize.py",
        ("token IDs → text streaming", 1),
        "server.py",
        ("FastAPI: /generate (full), /generate_stream (SSE)", 1),
    ], size=11)

    # Right column: phase mapping
    add_text(s, 7.2, 1.7, 5.5, 0.3, "What changes by phase",
             size=14, bold=True, color=TEAL)
    add_box(s, 7.2, 2.05, 5.5, 0.6, fill=LIGHT, border=TEAL,
            text="P0  one request at a time, full-only response", size=11, bold=True)
    add_box(s, 7.2, 2.7, 5.5, 0.6, fill=LIGHT, border=TEAL,
            text="P1  scheduler turns on continuous batching", size=11, bold=True)
    add_box(s, 7.2, 3.35, 5.5, 0.6, fill=WHITE, border=GREY,
            text="P2  no change (transparent KV swap)", size=11)
    add_box(s, 7.2, 4.0, 5.5, 0.6, fill=LIGHT, border=TEAL,
            text="P3  forward op awaits futures, no sync", size=11, bold=True)
    add_box(s, 7.2, 4.65, 5.5, 0.6, fill=WHITE, border=GREY,
            text="P4  no change", size=11)
    add_box(s, 7.2, 5.3, 5.5, 0.6, fill=WHITE, border=GREY,
            text="P5  no change", size=11)

    add_text(s, 7.2, 6.05, 5.5, 0.4,
             "Day 1: build against MOCK ForwardRunner.\n"
             "Develop & test fully WITHOUT an NPU.",
             size=11, italic=True, color=GREY)

    add_footer(s, n, total)


def slide_contract(prs, n, total):
    s = add_title_only(prs, "The contract — what crosses the boundary",
                       "Three artifacts, written together on day 1, then frozen")
    code = (
        "@dataclass\n"
        "class PrefillBatch:\n"
        "    request_ids: list[str]\n"
        "    input_ids: torch.Tensor       # (B, max_prompt_len) padded\n"
        "    attention_mask: torch.Tensor  # (B, max_prompt_len)\n\n"
        "@dataclass\n"
        "class DecodeBatch:\n"
        "    request_ids: list[str]\n"
        "    next_tokens: torch.Tensor     # (B,) sampled tokens from last step\n\n"
        "@dataclass\n"
        "class StepResult:\n"
        "    request_ids: list[str]\n"
        "    logits_future: asyncio.Future  # resolves to (B, vocab)\n\n"
        "class ForwardRunner:\n"
        "    def admit_request(self, request_id, prompt_len): ...\n"
        "    def release_request(self, request_id): ...\n"
        "    async def prefill_step(self, batch: PrefillBatch) -> StepResult: ...\n"
        "    async def decode_step(self, batch: DecodeBatch) -> StepResult: ...\n"
        "    def hbm_pressure(self) -> float: ...\n"
        "    def can_admit(self, prompt_len: int) -> bool: ..."
    )
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(12.3), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = code
    r.font.name = "Consolas"
    r.font.size = Pt(13)
    r.font.color.rgb = NAVY

    add_box(s, 0.5, 6.4, 12.3, 0.5, fill=ORANGE, border=ORANGE,
            text="Stable for v0–v3. Phase 5 may add a stage_id field; Component B unaffected.",
            size=12, bold=True, text_color=WHITE)
    add_footer(s, n, total)


def slide_phases(prs, n, total):
    s = add_title_only(prs, "Phase plan",
                       "Each phase ships working software; later phases additively optimize")
    rows = [
        ("P0  proof of life",          "single request, HF.generate, mock scheduler",
         "~600 LOC", "1 day"),
        ("P1  continuous batching",   "real prefill/decode tick, multi-tenant",
         "+400",     "2 days"),
        ("P2  DRAM KV",               "DramKVCache + contiguous per-request pool",
         "+250",     "1-2 days"),
        ("P3  streamed submission",   "futures, no per-step sync, instant admission",
         "+200",     "1 day"),
        ("P4  NPUGraph + FX offload", "decode-step graph capture + offload rewriter",
         "+400",     "2-3 days"),
        ("P5  PP + EP multi-NPU",     "partition_pipeline + EP rewrite + HCCL backend",
         "+500",     "3-5 days"),
    ]
    y = 1.7
    add_text(s, 0.5, y, 4.0, 0.4, "Phase",        size=13, bold=True, color=NAVY)
    add_text(s, 4.6, y, 5.5, 0.4, "Deliverable",  size=13, bold=True, color=NAVY)
    add_text(s, 10.2, y, 1.3, 0.4, "Code",         size=13, bold=True, color=NAVY)
    add_text(s, 11.6, y, 1.3, 0.4, "Effort",       size=13, bold=True, color=NAVY)
    y += 0.45
    for phase, deliv, loc, eff in rows:
        if phase.startswith("P5"):
            color = TEAL
        elif phase.startswith("P4"):
            color = ORANGE
        else:
            color = NAVY
        add_box(s, 0.5, y, 4.0, 0.6, fill=LIGHT, border=color,
                text=phase, size=12, bold=True)
        add_box(s, 4.6, y, 5.5, 0.6, fill=WHITE, border=GREY,
                text=deliv, size=11)
        add_box(s, 10.2, y, 1.3, 0.6, fill=WHITE, border=GREY,
                text=loc, size=11)
        add_box(s, 11.6, y, 1.3, 0.6, fill=WHITE, border=GREY,
                text=eff, size=11)
        y += 0.65
    add_text(s, 0.5, 6.0, 12.0, 0.4,
             "Total v0–v3 (working serving binary): ~1.5k LOC, 5–6 days for two people in parallel.",
             size=13, italic=True, color=GREY)
    add_footer(s, n, total)


def slide_timeline(prs, n, total):
    s = add_title_only(prs, "Timeline — two people, ~10 working days for v0–v3",
                       "Day 4 is the only required-sync day; everything else parallelizes")
    days = [
        ("Day 1", "JOINT — write contract files",     "JOINT"),
        ("Day 2", "ForwardRunner with HF.generate",   "Topology + admission + tokenize"),
        ("Day 3", "Manual decode loop, HF cache",     "Scheduler skeleton (FCFS)"),
        ("Day 4", "INTEGRATION DAY — both halves",    "INTEGRATION"),
        ("Day 5", "DramKVCache + kv_pool",            "Real continuous batching (P1)"),
        ("Day 6", "Test KV pool with N requests",     "Streaming output, FastAPI SSE"),
        ("Day 7", "Streamed submission (P3)",         "Backpressure, queue tuning"),
        ("Day 8", "NPUGraph capture for decode (P4)", "Stats endpoint, observability"),
        ("Day 9", "FX offload integration (P4)",      "Load testing, scheduler edge cases"),
        ("Day 10","Phase 4 integration test",         "Polish + docs"),
    ]
    add_text(s, 0.5, 1.7, 1.4, 0.35, "Day",      size=13, bold=True, color=NAVY)
    add_text(s, 2.0, 1.7, 5.5, 0.35, "Person 1 — Model Runtime",   size=13, bold=True, color=NAVY)
    add_text(s, 7.7, 1.7, 5.2, 0.35, "Person 2 — Serving Engine",  size=13, bold=True, color=TEAL)
    y = 2.1
    for day, p1, p2 in days:
        if "JOINT" in (p1, p2) or "INTEGRATION" in (p1, p2):
            fc = ORANGE
            tc = WHITE
        else:
            fc = LIGHT
            tc = BLACK
        add_box(s, 0.5, y, 1.4, 0.42, fill=fc, border=NAVY, text=day,
                size=11, bold=True, text_color=tc)
        add_box(s, 2.0, y, 5.5, 0.42, fill=fc if "JOINT" in p1 else LIGHT,
                border=NAVY, text=p1, size=10, bold="JOINT" in p1 or "INTEGRATION" in p1,
                text_color=tc if "JOINT" in p1 or "INTEGRATION" in p1 else BLACK)
        add_box(s, 7.7, y, 5.2, 0.42, fill=fc if "JOINT" in p2 or "INTEGRATION" in p2 else LIGHT,
                border=TEAL, text=p2, size=10,
                bold="JOINT" in p2 or "INTEGRATION" in p2,
                text_color=tc if "JOINT" in p2 or "INTEGRATION" in p2 else BLACK)
        y += 0.47
    add_footer(s, n, total)


def slide_day1(prs, n, total):
    s = add_title_only(prs, "Day 1 deliverables — write together, then freeze",
                       "Three files. Once compiled, Person 1 and Person 2 work in parallel.")

    add_box(s, 0.5, 1.7, 4.0, 1.2, fill=LIGHT, border=NAVY,
            text="config.py\n\nKnobs both sides agree on:\nmodel_id, dtype, device,\nmax_concurrent_requests,\nmax_seq_len, decode_batch_target",
            size=12, bold=True)

    add_box(s, 4.7, 1.7, 4.0, 1.2, fill=LIGHT, border=NAVY,
            text="state/request.py\n\nRequest, RequestState,\nSamplingParams,\nPrefillBatch / DecodeBatch,\nStepResult",
            size=12, bold=True)

    add_box(s, 8.9, 1.7, 4.0, 1.2, fill=LIGHT, border=NAVY,
            text="hf_integration/\nforward_runner.py\n\nABC + mock impl\n(sleep + random tokens)",
            size=12, bold=True)

    # Show the mock pattern
    code = (
        "# Used by Component B until Component A's real ForwardRunner ships.\n\n"
        "class MockForwardRunner(ForwardRunner):\n"
        "    async def prefill_step(self, batch: PrefillBatch) -> StepResult:\n"
        "        await asyncio.sleep(0.05 * len(batch.request_ids))    # fake compute\n"
        "        fake_logits = torch.randn(len(batch.request_ids), VOCAB)\n"
        "        fut = asyncio.get_running_loop().create_future()\n"
        "        fut.set_result(fake_logits)\n"
        "        return StepResult(batch.request_ids, fut)\n"
        "\n"
        "    async def decode_step(self, batch: DecodeBatch) -> StepResult:\n"
        "        await asyncio.sleep(0.005 * len(batch.request_ids))   # fake decode tick\n"
        "        ...\n"
        "\n"
        "    def hbm_pressure(self) -> float:    return 0.5\n"
        "    def can_admit(self, prompt_len) -> bool:    return True"
    )
    tb = s.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12.3), Inches(3.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = code
    r.font.name = "Consolas"
    r.font.size = Pt(12)
    r.font.color.rgb = NAVY

    add_text(s, 0.5, 6.7, 12.3, 0.4,
             "Person 2 builds the entire serving engine against this mock — no NPU needed for days 2-3.",
             size=12, italic=True, color=GREY)
    add_footer(s, n, total)


def slide_open_questions(prs, n, total):
    s = add_title_only(prs, "Open questions — decide before day 1",
                       "Four choices that shape the design")
    qs = [
        ("Output streaming", "SSE (simpler) vs WebSocket (richer)",
         "Recommend SSE for v0; FastAPI native"),
        ("Prefill chunking", "Chunked prefill (vLLM-style) vs full prefill per req",
         "Recommend full prefill for v0; chunk later"),
        ("Multi-NPU layout", "PP=4 × EP=4 vs PP=8 vs DP=8 (replicate)",
         "Recommend PP=4 × EP=4 for full-model fit"),
        ("Phase 4 scope",   "Offload + NPUGraph at MVP, or after serving binary lands?",
         "Recommend AFTER — get serving stable first"),
    ]
    y = 1.7
    add_text(s, 0.5, y, 3.0, 0.4, "Question",      size=13, bold=True, color=NAVY)
    add_text(s, 3.6, y, 5.0, 0.4, "Options",       size=13, bold=True, color=NAVY)
    add_text(s, 8.7, y, 4.2, 0.4, "Recommendation",size=13, bold=True, color=ORANGE)
    y += 0.45
    for q, opts, rec in qs:
        add_box(s, 0.5, y, 3.0, 0.85, fill=LIGHT, border=NAVY, text=q,
                size=12, bold=True)
        add_box(s, 3.6, y, 5.0, 0.85, fill=WHITE, border=GREY, text=opts, size=11)
        add_box(s, 8.7, y, 4.2, 0.85, fill=ORANGE, border=ORANGE, text=rec,
                size=11, bold=True, text_color=WHITE)
        y += 0.95
    add_footer(s, n, total)


def slide_summary(prs, n, total):
    s = add_title_only(prs, "Summary",
                       "What we're building, why, and how we split it")
    add_box(s, 0.5, 1.7, 12.3, 0.7, fill=NAVY, border=NAVY,
            text="A lite serving layer over HF transformers, optimized for DRAM-rich Ascend NPU box.",
            size=15, bold=True, text_color=WHITE)
    add_bullets(s, 0.5, 2.6, 12.3, 4.4, [
        "Two components, one contract: Model Runtime (NPU) ↔ Serving Engine (async)",
        ("ForwardRunner is the boundary. Mock impl on day 1 unblocks parallel work.", 1),
        "Continuous batching + streamed submission = primary throughput levers",
        ("Bigger batch is THE lever. No-sync forward path is the cheapest meaningful win.", 1),
        "DRAM-resident KV (contiguous per request, no paging)",
        ("2 TB DRAM removes vLLM's bin-packing constraint; simpler abstraction.", 1),
        "FX rewriter + NPUGraph for offload (phase 4, after serving works)",
        ("Capture once, replay every decode step; offload K/V projection weights by default.", 1),
        "PP + EP for multi-NPU; no tensor parallelism",
        ("PP=4 × EP=4 saturates 16 NPUs; avoids per-op rewriting cost of TP.", 1),
        "5–6 days for working serving binary (v0–v3); 8–12 days with offload + multi-NPU.",
    ], size=13)
    add_footer(s, n, total)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        slide_cover,        # cover doesn't need n/total
        slide_why,
        slide_hardware,
        slide_findings_helps,
        slide_findings_skip,
        slide_topology,
        slide_two_components,
        slide_component_a,
        slide_component_b,
        slide_contract,
        slide_phases,
        slide_timeline,
        slide_day1,
        slide_open_questions,
        slide_summary,
    ]
    total = len(builders)
    for i, fn in enumerate(builders, start=1):
        if fn is slide_cover:
            fn(prs)
        else:
            fn(prs, i, total)

    out = Path(__file__).resolve().parent.parent / "out" / "hfc_design.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"wrote {out} ({out.stat().st_size / 1024:.1f} KB, {total} slides)")


if __name__ == "__main__":
    main()
